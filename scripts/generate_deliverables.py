"""
Génère les livrables rédigés au format Office (NON commités dans le dépôt) :
  livrables/rapport_projet3.docx
  livrables/plan_soutenance_projet3.docx
  livrables/antiseche_projet3.docx
  livrables/dossier_audit_projet3.docx
  livrables/presentation_projet3.pptx

Les chiffres clés sont lus dynamiquement depuis evaluation/*.md quand ils
existent (sinon valeurs de repli). Captures réelles : screenshots/*.png ;
figure d'éval : evaluation/figures/*.png.

Usage : python scripts/generate_deliverables.py
"""

import os
import re
from pathlib import Path

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "livrables"
OUT.mkdir(exist_ok=True)
SHOTS = ROOT / "screenshots"
FIGS = ROOT / "evaluation" / "figures"

AUTHORS = "BOUCHER Anthony & DEROUICHE Youcef"
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
PACCENT = PRGB(0x0E, 0xA5, 0xE9)

# --- Chiffres clés (repli ; surchargés par lecture des résultats si dispo) ---
NUM = {
    "ndcg_sem": "0.333", "ndcg_bug": "0.326", "ndcg_cal": "0.506",
    "mrr_sem": "0.593", "mrr_cal": "0.778", "gain_ndcg": "+52",
    "tfidf_ndcg": "0.031", "sbert_ndcg": "0.333", "tfidf_gain": "+961",
    "w_naive": "0.489", "w_chosen": "0.506",
    "model": "gemini-2.5-flash",
}


def _read(path):
    p = ROOT / path
    return p.read_text(encoding="utf-8") if p.exists() else ""


def _load_numbers():
    """Met à jour NUM depuis les fichiers de résultats générés."""
    res = _read("evaluation/RESULTS.md")
    m = re.search(r"de \*\*([\d.]+)\*\* à \*\*([\d.]+)\*\* \(\+?([\d.]+)\s*%", res)
    if m:
        NUM["ndcg_sem"], NUM["ndcg_cal"] = m.group(1), m.group(2)
        NUM["gain_ndcg"] = "+" + m.group(3).split(".")[0]
    base = _read("evaluation/baseline_results.md")
    mb = re.search(r"nDCG@5 de \*\*\+?([\d.]+)\s*%", base)
    if mb:
        NUM["tfidf_gain"] = "+" + mb.group(1).split(".")[0]


# ---------------------------------------------------------------- DOCX helpers
def h1(doc, t): return doc.add_heading(t, level=1)
def h2(doc, t): return doc.add_heading(t, level=2)


def para(doc, text, bold=False, italic=False):
    p = doc.add_paragraph(); r = p.add_run(text); r.bold = bold; r.italic = italic
    return p


def bullet(doc, text): doc.add_paragraph(text, style="List Bullet")
def numbered(doc, text): doc.add_paragraph(text, style="List Number")


def table(doc, headers, rows):
    t = doc.add_table(rows=1, cols=len(headers)); t.style = "Light Grid Accent 1"
    for i, h in enumerate(headers):
        run = t.rows[0].cells[i].paragraphs[0].add_run(h); run.bold = True
    for row in rows:
        cells = t.add_row().cells
        for i, v in enumerate(row):
            cells[i].text = str(v)
    return t


def title_block(doc, title, subtitle):
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title); r.bold = True; r.font.size = Pt(24); r.font.color.rgb = ACCENT
    p2 = doc.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.add_run(subtitle).font.size = Pt(12)
    p3 = doc.add_paragraph(); p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.add_run(AUTHORS + "  —  RNCP40875, Bloc 2 (C5.1 → C5.3)").italic = True
    doc.add_paragraph()


# ======================================================= 1. RAPPORT
def build_rapport():
    doc = Document()
    title_block(doc, "Rapport — Projet 3 : IA Générative",
                "AISCA-Cinema : agent de recommandation cinématographique (RAG)")

    h1(doc, "1. Besoin métier (C5.1)")
    para(doc, "Face au volume de films disponibles, choisir quoi regarder est coûteux. "
              "Les moteurs par mots-clés ne comprennent pas une envie exprimée en langage "
              "naturel et n'expliquent pas leurs suggestions. AISCA-Cinema répond au besoin "
              "de découverte personnalisée : l'utilisateur décrit son envie avec ses mots, "
              "et l'agent recommande des films pertinents ET justifie son choix, à partir "
              "d'un référentiel maîtrisé (pas d'invention).")
    para(doc, "Pourquoi la GenAI : compréhension sémantique de la demande (SBERT) + "
              "restitution personnalisée en langage naturel (LLM), hors de portée d'un "
              "simple filtre par genres.")

    h1(doc, "2. Choix techniques")
    para(doc, "Approche retenue : un RAG (Retrieval-Augmented Generation) — récupération "
              "sémantique de films pertinents puis génération de synthèses ancrées sur ces films.")
    table(doc, ["Décision", "Choix", "Justification"], [
        ["Paradigme", "RAG", "corpus évolutif sans réentraînement ; génération traçable ; coût maîtrisé"],
        ["vs Fine-tuning", "écarté", "coût/temps, données insuffisantes, perte de traçabilité"],
        ["vs Prompting seul", "écarté", "risque d'hallucination sans ancrage corpus"],
        ["Embeddings", "SBERT multilingue MiniLM-L12-v2", "requête FR / corpus EN, léger en CPU"],
        ["Génération", f"Google Gemini ({NUM['model']})", "qualité FR, free tier, latence faible"],
        ["Coûts API", "cache + retry/backoff + 2 appels/session", "free tier respecté, résilience"],
    ])
    para(doc, "Pourquoi SBERT plutôt qu'un index lexical : une baseline TF-IDF (mots-clés) "
              f"obtient un nDCG@5 de {NUM['tfidf_ndcg']} contre {NUM['sbert_ndcg']} pour SBERT "
              f"(soit {NUM['tfidf_gain']} %). L'écart est massif car les requêtes sont en "
              "français et les descriptions en anglais : seul un modèle multilingue fait le "
              "pont sémantique — ce qu'un index de mots-clés ne peut pas.")

    h1(doc, "3. Réalisation")
    bullet(doc, "questionnaire.py : acquisition et structuration des préférences.")
    bullet(doc, "nlp_engine.py : retrieval SBERT + cache des embeddings (mémoire ET disque) + index FAISS optionnel.")
    bullet(doc, "scoring.py : ranking pondéré 0.50 sémantique + 0.40 genre + 0.10 mood (poids calibrés, cf. §4).")
    bullet(doc, "genai_integration.py : génération Gemini, GenerationConfig ajustable, retry/backoff (quota 429), "
                "safety_settings, garde-fou anti-prompt-injection, suivi tokens/coût/latence, LLM-as-judge.")
    bullet(doc, "app.py : UI Streamlit accessible, traçabilité des sources, mode dégradé sans clé API.")
    h2(doc, "Correctifs majeurs apportés lors de la finalisation")
    numbered(doc, "Sécurité : retrait d'une clé API réelle commitée dans .env.example (rotation effectuée).")
    numbered(doc, "Bug fonctionnel (C5.2) : composant genre neutralisé (libellés FR vs colonne Genre "
                  "anglaise). Corrigé via la colonne Categorie (FR) + matching insensible aux accents.")
    numbered(doc, "Évaluation (C5.3) : cadre d'évaluation complet créé (métriques retrieval, validation "
                  "croisée des poids, LLM-as-judge, comparaison de paramètres).")
    numbered(doc, "Industrialisation : Dockerfile, persistance des embeddings, suivi des coûts, CI GitHub Actions.")

    h1(doc, "4. Preuve (mesures réelles)")
    h2(doc, "4.1 Qualité du RAG — comparaison avant / après")
    para(doc, "Jeu de cas : 15 requêtes annotées (vérité terrain). Métriques @5 :")
    table(doc, ["Configuration", "MRR", "nDCG@5"], [
        ["A. Sémantique seule (baseline)", NUM["mrr_sem"], NUM["ndcg_sem"]],
        ["B. Pondéré — genre buggé (colonne EN)", "0.533", NUM["ndcg_bug"]],
        ["C. Pondéré — corrigé + calibré 50/40/10", NUM["mrr_cal"], NUM["ndcg_cal"]],
    ])
    para(doc, f"Le correctif (A → C) améliore le nDCG@5 de {NUM['gain_ndcg']} % et le MRR de "
              f"{NUM['mrr_sem']} à {NUM['mrr_cal']}. La config B (genre buggé) reste collée au "
              "sémantique seul — preuve chiffrée que le composant genre était inopérant.")
    if (FIGS / "rag_comparison.png").exists():
        doc.add_picture(str(FIGS / "rag_comparison.png"), width=Inches(5.6))
    h2(doc, "4.2 Calibrage des poids (validation croisée)")
    para(doc, "Les poids ne sont pas fixés à la main : un grid-search sur le jeu annoté "
              f"montre que 0.50/0.40/0.10 (nDCG@5 {NUM['w_chosen']}) bat le réglage naïf "
              f"0.50/0.30/0.20 ({NUM['w_naive']}). On garde le sémantique dominant (robuste "
              "sur les requêtes libres) et le mood comme départage léger. Script : "
              "evaluation/tune_weights.py.")
    h2(doc, "4.3 Qualité de la génération (LLM-as-judge + paramètres)")
    para(doc, "Un LLM-as-judge (appel déterministe, JSON strict) note chaque réponse de 1 à 5 "
              "(pertinence, exactitude, complétude, absence d'hallucination, ton). On compare "
              "deux réglages — Factuelle (température 0.2) vs Créative (0.9) — et on retient le "
              "meilleur compromis qualité/fiabilité. Scripts : evaluation/compare_params.py "
              "(résultats dans evaluation/params_results.md).")
    h2(doc, "4.4 Tests & CI")
    para(doc, "29 tests unitaires (scoring, métriques, cache, anti-injection) ; CI GitHub "
              "Actions exécutant les tests à chaque push (garde-fou de non-régression).")

    h1(doc, "5. Résultat")
    para(doc, "Application fonctionnelle et accessible, RAG réellement personnalisé (correctif "
              "prouvé), génération ancrée/traçable et résiliente, évaluation chiffrée avec "
              "comparaison avant/après — C5.1, C5.2 et C5.3 couvertes.")

    h1(doc, "6. Limites, biais et risques")
    bullet(doc, "Corpus : 260 films classiques (IMDb), descriptions EN → biais culturel / couverture limitée.")
    bullet(doc, "Hallucination : atténuée (prompts + sources + garde-fou) mais résiduelle hors top 3.")
    bullet(doc, "Dépendance API + quotas free tier (limites/jour et /minute) → mode dégradé + retry/backoff.")
    bullet(doc, "RGPD : stockage local des réponses libres à encadrer (consentement, rétention).")
    bullet(doc, "Évaluation : vérité terrain construite par les auteurs (subjectivité) ; à élargir.")

    h1(doc, "7. Industrialisation")
    bullet(doc, "Index vectoriel (FAISS, déjà intégré en option) + embeddings persistés (fait).")
    bullet(doc, "Docker (fait) + CI exécutant tests (fait) ; étendre la CI à l'évaluation.")
    bullet(doc, "Monitoring coûts/latence (instrumenté) + supervision du taux d'hallucination + A/B testing des paramètres.")

    h1(doc, "8. Compétences démontrées (mapping RNCP)")
    table(doc, ["Compétence", "Réalisation", "Preuve"], [
        ["C5.1 Cas d'usage", "Besoin métier formalisé, justification RAG, gouvernance (clé, safety, RGPD)", "§1, §2, §6"],
        ["C5.2 Solution", "RAG fonctionnel, accessible, traçable, résilient ; correctif scoring", "app, §3, tests"],
        ["C5.3 Évaluation", "Métriques + validation croisée poids + LLM-as-judge + comparaison params", f"evaluation/, {NUM['gain_ndcg']}% nDCG"],
    ])

    h1(doc, "9. Ma contribution individuelle (à compléter par chaque membre)")
    para(doc, "À personnaliser pour la défense individualisée :", italic=True)
    bullet(doc, "[Prénom] : moteur de retrieval SBERT, système de scoring, correction du bug genre/mood, "
                "calibrage des poids, tests. Choix défendus : SBERT multilingue, pondération calibrée, colonne Categorie.")
    bullet(doc, "[Prénom] : intégration Gemini (prompts, GenerationConfig, retry, safety, anti-injection, coûts), "
                "cadre d'évaluation (métriques, LLM-as-judge, compare_params), UI et traçabilité, Docker/CI.")

    doc.save(OUT / "rapport_projet3.docx")
    print("OK rapport_projet3.docx")


# ======================================================= 2. PLAN DE SOUTENANCE
def build_plan():
    doc = Document()
    title_block(doc, "Plan de soutenance — Projet 3 (IA Générative)",
                "Déroulé slide par slide · qui dit quoi (mot à mot) · transitions · timing")
    para(doc, "Format global : 3 projets en 30 min au total (~10 min/projet) puis 20 min de "
              "questions du jury. Objectif de ce projet : ~10 min, 9 slides.", italic=True)
    para(doc, "Règle d'or : chacun défend ce qu'il présente. Anthony porte le fil « moteur & "
              "données » (retrieval, scoring, correctif, évaluation) ; Youcef porte le fil "
              "« IA générative & produit » (architecture, choix, LLM, robustesse, industrialisation). "
              "Chacun doit toutefois pouvoir répondre sur tout (cf. antisèche).", italic=True)

    slides = [
        ("Slide 1 — Titre", "Titre AISCA-Cinema, sous-titre, noms, RNCP.", "—", "Anthony", "0:30",
         "« Bonjour, nous sommes Anthony Boucher et Youcef Derouiche. Nous présentons AISCA-Cinema, "
         "notre projet d'IA générative pour le bloc 2 : un agent de recommandation de films fondé "
         "sur une architecture RAG, c'est-à-dire génération augmentée par récupération. »",
         "→ Anthony enchaîne sur le besoin métier."),
        ("Slide 2 — Besoin métier (C5.1)", "Problème, cible, valeur, pourquoi la GenAI.", "—",
         "Anthony", "1:00",
         "« Le problème est simple : il y a trop de films, et choisir coûte du temps. Les moteurs "
         "classiques filtrent par mots-clés ou par genre, mais ils ne comprennent pas une envie "
         "formulée en langage naturel — par exemple « un film contemplatif sur la mémoire » — et "
         "surtout ils n'expliquent pas leurs suggestions. Notre cible, c'est un spectateur qui sait "
         "décrire une envie mais pas un titre. La valeur qu'on apporte : recommander ET justifier, "
         "à partir d'un référentiel maîtrisé. Et pourquoi de l'IA générative ici ? Parce que la "
         "valeur tient à deux choses qu'un simple filtre ne sait pas faire : comprendre le sens de "
         "la demande, et restituer une explication personnalisée. C'est notre cas d'usage, la "
         "compétence C5.1. »",
         "→ Anthony passe la main à Youcef : « Youcef va détailler l'architecture. »"),
        ("Slide 3 — Architecture RAG", "Schéma pipeline (4 cartes) + capture questionnaire.",
         "01_questionnaire.png", "Youcef", "1:30",
         "« Voici le pipeline. L'utilisateur remplit un questionnaire — une description libre plus "
         "des curseurs de préférences. Ce texte part dans SBERT, un modèle qui transforme la phrase "
         "en vecteur numérique ; on encode de la même façon nos 260 films, et on mesure la proximité "
         "par similarité cosinus : c'est l'étape de récupération, le Retrieval. Ensuite on reclasse "
         "les films par un score pondéré — 50 % sémantique, 40 % genre, 10 % ambiance. Enfin, Gemini, "
         "le LLM de Google, rédige un profil et un plan de découverte, mais uniquement à partir des "
         "films récupérés : c'est l'étape de génération, la Generation. Retrieval plus Generation, "
         "c'est exactement ce qu'on appelle un RAG. »",
         "→ Youcef enchaîne sur la justification des choix."),
        ("Slide 4 — Choix techniques (C5.2)", "Tableau RAG vs alternatives + baseline TF-IDF.",
         "—", "Youcef", "1:15",
         "« Pourquoi un RAG et pas autre chose ? Le fine-tuning aurait demandé beaucoup de données "
         "et serait devenu opaque ; le prompting seul aurait halluciné faute d'ancrage. Le RAG, lui, "
         "laisse le corpus évoluer sans réentraînement, reste traçable — on sait quels films ont servi "
         "— et coûte peu. Deuxième choix : pourquoi des embeddings plutôt qu'une recherche par "
         "mots-clés ? On l'a mesuré : une baseline TF-IDF s'effondre face à SBERT, plus de 900 % "
         "d'écart de nDCG. La raison est parlante : nos requêtes sont en français et les films décrits "
         "en anglais — seul un modèle multilingue fait le pont sémantique. C'est notre justification "
         "C5.2. »",
         "→ Youcef repasse la main à Anthony pour la démo."),
        ("Slide 5 — Démo Top 3", "Capture des recommandations réelles + composantes du score.",
         "02_top3.png", "Anthony", "1:30",
         "« Voici une vraie sortie de l'application, pour une requête de science-fiction "
         "philosophique. On obtient le Top 3, et surtout, pour chaque film, la décomposition du "
         "score : la part sémantique, la part genre, la part ambiance. La recommandation n'est pas "
         "une boîte noire : elle est expliquée et tracée jusqu'aux films sources. »",
         "→ Anthony bascule sur l'onglet visualisations."),
        ("Slide 6 — Démo Visualisations", "Radars de préférences + barres de score 50/40/10.",
         "03_visualisations.png", "Youcef", "0:45",
         "« L'interface est pensée pour un non-technicien : deux radars résument le profil par genre "
         "et par ambiance, et un graphique montre la décomposition pondérée du score de chaque film "
         "recommandé. C'est le volet accessibilité de la compétence C5.2. »",
         "→ Youcef passe à Anthony pour l'évaluation, le cœur de C5.3."),
        ("Slide 7 — Évaluation (C5.3)", "Figure avant/après + cartes-chiffres + poids + juge.",
         "rag_comparison.png", "Anthony", "2:00",
         "« C'est le cœur de la compétence C5.3, l'évaluation. On a construit un jeu de 15 requêtes "
         "annotées, avec une vérité terrain : pour chaque requête, la liste des films réellement "
         "pertinents. On mesure quatre métriques standard de recherche d'information : Precision, "
         "Recall, MRR et nDCG. Résultat : notre correctif fait passer le nDCG de 0,33 à 0,51, plus "
         "de 50 % de gain, et le MRR de 0,59 à 0,78. La barre orange, au milieu, c'est la version "
         "où le genre était buggé : elle reste collée au sémantique seul, ce qui prouve chiffres à "
         "l'appui que le composant genre était mort avant notre correction. Point important pour le "
         "jury : les poids 50/40/10 ne sont pas choisis au doigt mouillé — ils sont validés par un "
         "grid-search, une recherche exhaustive sur le jeu annoté. Et pour la qualité des textes "
         "générés par le LLM, on utilise un LLM-as-judge : un second appel du modèle note chaque "
         "réponse sur des critères comme la pertinence et l'absence d'hallucination, ce qui nous "
         "permet de comparer objectivement deux réglages de température. »",
         "→ Anthony passe la main à Youcef pour la robustesse et les limites."),
        ("Slide 8 — Robustesse · Sécurité · Limites", "Résilience, gouvernance, observabilité, limites.",
         "—", "Youcef", "1:00",
         "« Côté production, on a soigné trois choses. La résilience : des relances automatiques "
         "avec backoff exponentiel quand l'API renvoie un dépassement de quota, plus un mode dégradé "
         "qui garde le cœur fonctionnel même sans clé. La gouvernance : des safety_settings et un "
         "garde-fou contre l'injection de prompt sur le texte libre de l'utilisateur. Et "
         "l'observabilité : on suit les tokens, le coût estimé et la latence à chaque appel. Nos "
         "limites, qu'on assume : un corpus de 260 films classiques en anglais, une hallucination "
         "résiduelle possible hors du top 3, et les quotas du free tier. »",
         "→ Youcef conclut sur l'industrialisation."),
        ("Slide 9 — Industrialisation & conclusion", "Pistes d'industrialisation + récap compétences.",
         "—", "Youcef puis Anthony", "0:30",
         "Youcef : « Pour industrialiser : un index vectoriel FAISS et des embeddings persistés — "
         "déjà intégrés —, Docker et une CI qui lance les tests à chaque commit, du monitoring des "
         "coûts, et de l'A/B testing des paramètres. » "
         "Anthony (clôture) : « En résumé : un cas d'usage justifié et gouverné pour C5.1, une "
         "solution RAG fonctionnelle et résiliente pour C5.2, et une évaluation chiffrée avec "
         "optimisation prouvée pour C5.3. Merci, nous sommes à votre disposition pour vos questions. »",
         "→ Fin de la présentation, place aux questions."),
    ]
    for titre, contenu, img, qui, timing, script, transition in slides:
        h1(doc, titre)
        table(doc, ["Élément", "Détail"], [["À l'écran", contenu], ["Image", img],
                                            ["Intervenant", qui], ["Durée", timing]])
        para(doc, "Script (mot à mot) :", bold=True)
        para(doc, script, italic=True)
        para(doc, "Transition : ", bold=True)
        para(doc, transition)
        doc.add_paragraph()

    h1(doc, "Récapitulatif timing")
    table(doc, ["Bloc", "Intervenant", "Durée cumulée"], [
        ["Slides 1-2 (intro + besoin)", "Anthony", "~1:30"],
        ["Slides 3-4 (archi + choix)", "Youcef", "~2:45"],
        ["Slide 5 (démo Top 3)", "Anthony", "~1:30"],
        ["Slide 6 (visualisations)", "Youcef", "~0:45"],
        ["Slide 7 (évaluation)", "Anthony", "~2:00"],
        ["Slides 8-9 (robustesse + conclusion)", "Youcef (+ Anthony clôture)", "~1:30"],
        ["TOTAL", "—", "~10:00"],
    ])

    h1(doc, "Répartition des questions du jury (qui répond en premier)")
    para(doc, "Le premier nommé répond ; l'autre complète. Les deux doivent connaître l'antisèche.")
    table(doc, ["Type de question", "Répond en premier"], [
        ["Besoin métier, cas d'usage, valeur, RGPD/gouvernance", "Anthony"],
        ["Retrieval, SBERT, embeddings, similarité, scoring, correctif genre", "Anthony"],
        ["Évaluation, métriques, validation croisée des poids, vérité terrain", "Anthony"],
        ["Architecture RAG, choix RAG vs fine-tuning/prompting", "Youcef"],
        ["LLM, Gemini, température/top-p, hallucination, LLM-as-judge", "Youcef"],
        ["Résilience, quotas, sécurité (injection, safety), coûts", "Youcef"],
        ["Industrialisation, FAISS, Docker, CI, monitoring", "Youcef"],
    ])

    h1(doc, "Conseils de présentation")
    for c in ["Relier explicitement chaque partie à une compétence (« ceci démontre Cx.y car… »).",
              "Toujours donner un chiffre quand on parle de résultat (nDCG, MRR, gain %).",
              "Assumer les limites : le jury valorise le recul critique.",
              "En cas de trou, se raccrocher à l'antisèche (section correspondante).",
              "Répéter au moins une fois à blanc avec un chrono pour tenir les 10 minutes."]:
        bullet(doc, c)

    doc.save(OUT / "plan_soutenance_projet3.docx")
    print("OK plan_soutenance_projet3.docx")


# ======================================================= 3. ANTISÈCHE
def build_antiseche():
    doc = Document()
    title_block(doc, "Antisèche technique — Projet 3 (IA Générative)",
                "Chaque terme expliqué au niveau ingénierie : définition · mécanisme · rôle · forces/limites")
    para(doc, "Objectif : pouvoir expliquer et défendre chaque brique du projet. Pour chaque concept : "
              "ce que c'est, comment ça marche techniquement, à quoi il sert chez nous, et ses "
              "forces/limites.", italic=True)

    def concept(nom, definition, mecanisme, role, ff):
        h2(doc, nom)
        para(doc, "Définition : ", bold=True); para(doc, definition)
        para(doc, "Comment ça marche : ", bold=True); para(doc, mecanisme)
        para(doc, "Rôle dans le projet : ", bold=True); para(doc, role)
        para(doc, "Forces / limites : ", bold=True); para(doc, ff)

    # ------------------------------------------------ A. NLP / RETRIEVAL
    h1(doc, "A. Traitement du langage & récupération (retrieval)")

    concept("NLP (Traitement Automatique du Langage)",
            "Ensemble de techniques pour faire manipuler du langage humain par une machine.",
            "On transforme du texte en représentations numériques exploitables (vecteurs), sur "
            "lesquelles on calcule des similarités, des classifications, etc.",
            "Toute la partie « comprendre la requête » et « comparer aux films » relève du NLP.",
            "Force : capture le sens au-delà des mots exacts. Limite : dépend de la qualité du "
            "modèle et des données.")

    concept("Transformer & mécanisme d'attention",
            "Architecture de réseau de neurones (2017) à la base des LLM et de BERT.",
            "Le mécanisme d'« attention » pondère l'importance de chaque mot par rapport aux autres "
            "dans la phrase, ce qui permet de capturer le contexte (un mot n'a pas le même sens "
            "selon son entourage). Les mots sont d'abord découpés en tokens puis projetés en vecteurs.",
            "SBERT et Gemini reposent tous deux sur des Transformers.",
            "Force : contexte long, parallélisable. Limite : coûteux en calcul, gourmand en données.")

    concept("Tokenisation",
            "Découpage d'un texte en unités élémentaires (tokens : mots, sous-mots, ponctuation).",
            "Un tokenizer applique un vocabulaire appris (ex. sous-mots BPE/WordPiece) : « cinéphile » "
            "peut devenir « ciné » + « phile ». Le modèle ne voit que des identifiants de tokens. "
            "La facturation des LLM et les limites de contexte se comptent en tokens.",
            "Détermine la longueur/coût des appels Gemini ; nos réponses sont bornées en tokens.",
            "Force : gère les mots inconnus via les sous-mots. Limite : un même texte peut coûter "
            "plus de tokens selon la langue.")

    concept("Embeddings (plongements vectoriels)",
            "Représentation d'un texte par un vecteur de nombres réels (ici ~384 dimensions).",
            "Le modèle place les textes dans un espace vectoriel tel que deux textes de sens proche "
            "ont des vecteurs proches. Chaque dimension n'est pas interprétable seule ; c'est la "
            "géométrie globale (distances/angles) qui porte le sens.",
            "On encode la requête utilisateur et chaque film ; la comparaison se fait sur ces vecteurs.",
            "Force : compare le SENS, multilingue. Limite : « boîte noire », sensible au modèle choisi.")

    concept("BERT vs Sentence-BERT (SBERT) — modèle MiniLM-L12-v2 multilingue",
            "BERT produit des vecteurs par token ; Sentence-BERT est optimisé pour produire un seul "
            "vecteur par PHRASE, directement comparable.",
            "SBERT ajoute une couche de « pooling » sur BERT et est entraîné par paires de phrases "
            "(objectif contrastif) pour que la distance vectorielle reflète la similarité sémantique. "
            "La version multilingue partage un espace commun entre langues (FR et EN).",
            "C'est notre moteur de récupération (src/nlp_engine.py) ; il gère nos requêtes FR sur "
            "des descriptions EN.",
            "Force : rapide, léger (tourne en CPU), multilingue. Limite : modèle généraliste, non "
            "spécialisé cinéma.")

    concept("Normalisation L2 & similarité cosinus",
            "Mesure de proximité entre deux vecteurs, insensible à leur longueur.",
            "Le cosinus de l'angle entre deux vecteurs vaut 1 s'ils pointent dans la même direction "
            "(très similaires), 0 s'ils sont orthogonaux (sans rapport). La normalisation L2 ramène "
            "les vecteurs à une norme 1, ce qui rend le produit scalaire égal au cosinus (utilisé "
            "par FAISS IndexFlatIP).",
            "Sert à classer les films par proximité avec la requête (calculate_similarity).",
            "Force : standard, robuste à la longueur du texte. Limite : ne capte que ce que "
            "l'embedding encode.")

    concept("TF-IDF (baseline lexicale)",
            "Représentation classique d'un texte par la fréquence pondérée de ses mots.",
            "TF = fréquence du terme dans le document ; IDF = rareté du terme dans le corpus. Un mot "
            "fréquent dans un doc mais rare globalement pèse fort. On compare ensuite par cosinus. "
            "C'est purement lexical : aucun sens, seulement des correspondances de mots.",
            "Notre baseline de comparaison : elle prouve l'apport de SBERT (+961 % de nDCG) car "
            "elle échoue sur le pont FR→EN.",
            "Force : simple, rapide, interprétable. Limite : aucun sens, ne franchit pas la barrière "
            "de langue ni les synonymes.")

    concept("Retrieval vs Ranking",
            "Deux étapes : récupérer un ensemble de candidats, puis les ordonner finement.",
            "Le retrieval (SBERT + cosinus) sélectionne les films proches de la requête ; le ranking "
            "réordonne ce sous-ensemble avec un score plus riche (sémantique + préférences). On "
            "sépare les deux pour la performance et la clarté.",
            "Retrieval = nlp_engine ; Ranking = scoring.py.",
            "Force : modulaire, chaque étape optimisable. Limite : une erreur de retrieval ne peut "
            "pas être rattrapée par le ranking.")

    concept("Score pondéré & normalisation Likert",
            "Combinaison linéaire de trois signaux : 0.50 sémantique + 0.40 genre + 0.10 ambiance.",
            "Les curseurs Likert (1 à 5) sont normalisés en [0,1] (division par 5). Les composantes "
            "genre/mood comparent la catégorie française du film aux préférences (matching insensible "
            "aux accents). Le score final est borné dans [0,1].",
            "Cœur du reclassement (scoring.py) ; c'est là qu'était le bug (colonne EN au lieu de FR).",
            "Force : transparent, explicable, ajustable. Limite : pondération linéaire, ne capture "
            "pas d'interactions complexes.")

    concept("FAISS & index vectoriel",
            "Bibliothèque de recherche de plus proches voisins dans un grand ensemble de vecteurs.",
            "Au lieu de comparer la requête à tous les films un par un (coûteux à grande échelle), "
            "FAISS construit un index (ici IndexFlatIP, produit scalaire sur vecteurs normalisés = "
            "cosinus) qui accélère la recherche. Les embeddings sont aussi persistés sur disque "
            "(.npy) pour éviter de les recalculer au démarrage.",
            "Intégré en option (nlp_engine.build_faiss_index) pour la scalabilité / industrialisation.",
            "Force : latence et passage à l'échelle. Limite : index à reconstruire si le corpus change.")

    # ------------------------------------------------ B. LLM / GÉNÉRATION
    h1(doc, "B. IA générative (LLM)")

    concept("LLM & génération autoregressive",
            "Large Language Model : modèle génératif entraîné à prédire le mot (token) suivant.",
            "Le modèle génère du texte token par token : à chaque étape il calcule une distribution "
            "de probabilité sur le vocabulaire et en tire le token suivant, qu'il réinjecte. Répété, "
            "cela produit une réponse cohérente.",
            "Gemini rédige le profil cinéphile et le plan de découverte.",
            "Force : qualité rédactionnelle, polyvalence. Limite : peut halluciner ; dépend d'une API.")

    concept("RAG (Retrieval-Augmented Generation)",
            "Architecture qui couple une RÉCUPÉRATION d'information et une GÉNÉRATION par LLM.",
            "On récupère d'abord des documents pertinents (nos films), qu'on injecte dans le prompt "
            "comme contexte ; le LLM répond en s'appuyant sur ce contexte plutôt que sur sa seule "
            "mémoire. Cela réduit les hallucinations et rend la réponse traçable.",
            "Toute notre solution : SBERT récupère, Gemini génère à partir des films récupérés.",
            "Force : ancrage factuel, corpus modifiable sans réentraînement, traçable. Limite : "
            "qualité plafonnée par le retrieval.")

    concept("Gemini 2.5 Flash & tokens de raisonnement",
            "LLM de Google ; la variante Flash est rapide et économique.",
            "C'est un modèle « à raisonnement » : il peut consommer une partie de son budget de "
            "sortie en réflexion interne avant de répondre. D'où un max_output_tokens élevé (3072) "
            "pour que la réponse visible ne soit pas tronquée.",
            "Modèle de génération du projet ; choisi car le free tier de gemini-2.0 renvoyait un "
            "quota nul.",
            "Force : bon rapport qualité/coût, multilingue. Limite : quotas free tier stricts "
            "(par jour ET par minute).")

    concept("Paramètres d'échantillonnage : température, top-p, top-k, max_output_tokens",
            "Réglages qui contrôlent l'aléa et la longueur de la génération.",
            "Température : aplatit (haute) ou accentue (basse) la distribution de probabilité — 0 = "
            "quasi déterministe/factuel, ~1 = créatif/varié. Top-k : ne considère que les k tokens "
            "les plus probables. Top-p (nucleus) : ne considère que les tokens dont la probabilité "
            "cumulée atteint p. max_output_tokens : longueur maximale.",
            "Exposés via GenerationConfig ; on compare Factuelle (0.2) vs Créative (0.9) dans "
            "compare_params.py.",
            "Force : arbitrer fiabilité vs richesse. Limite : température élevée = plus "
            "d'hallucinations.")

    concept("Prompt engineering & garde-fous anti-hallucination",
            "Conception des instructions envoyées au LLM pour obtenir un comportement fiable.",
            "On structure le prompt (rôle, données, tâche, contraintes de format et de longueur) et "
            "on ajoute des consignes explicites : « n'invente aucun fait sur les films listés », "
            "« appuie-toi uniquement sur le contexte fourni ». Le texte utilisateur est encadré par "
            "des balises pour ne pas être interprété comme une instruction.",
            "Utilisé dans generate_cinephile_profile / generate_discovery_plan.",
            "Force : réduit fortement les dérives. Limite : ne garantit pas 0 hallucination.")

    concept("Hallucination",
            "Production par le LLM d'un contenu plausible mais faux (film inexistant, fait inventé).",
            "Provient de la nature probabiliste du modèle : il complète ce qui « sonne juste ». On "
            "l'atténue par l'ancrage RAG, les consignes de prompt, une température basse, et on la "
            "mesure (titres cités hors corpus + critère « non-hallucination » du juge).",
            "Risque central d'un projet GenAI ; identifié comme limite assumée.",
            "Force (de notre approche) : plusieurs garde-fous cumulés. Limite : résiduelle hors top 3.")

    concept("LLM-as-judge",
            "Utiliser un LLM comme évaluateur automatique des réponses d'un autre (ou du même) LLM.",
            "Un appel dédié, en température 0 et sortie JSON stricte, note la réponse de 1 à 5 sur "
            "des critères (pertinence, exactitude, complétude, non-hallucination, ton). On extrait le "
            "JSON et on agrège. On utilise un modèle léger pour éviter la troncature par raisonnement.",
            "judge_response() ; alimente la comparaison de paramètres (qualité avant/après température).",
            "Force : évaluation qualité automatisée et reproductible. Limite : le juge a ses propres "
            "biais ; à recouper avec des notes humaines.")

    # ------------------------------------------------ C. ÉVALUATION
    h1(doc, "C. Évaluation (C5.3)")

    concept("Vérité terrain & jeu de cas annoté",
            "Ensemble de requêtes pour lesquelles on a défini à la main les films pertinents.",
            "On a rédigé 15 requêtes réalistes et, pour chacune, la liste des films du corpus jugés "
            "pertinents (evaluation/test_queries.json). C'est la référence contre laquelle on mesure "
            "le système.",
            "Base de toutes nos métriques et du calibrage des poids.",
            "Force : permet une évaluation objective. Limite : annotation par les auteurs "
            "(subjectivité), 15 requêtes = échantillon modeste.")

    concept("Precision@k, Recall@k, MRR, nDCG",
            "Métriques standard de recherche d'information.",
            "Precision@k : proportion de pertinents parmi les k premiers. Recall@k : proportion des "
            "pertinents retrouvés dans les k premiers. MRR : moyenne de 1/rang du premier résultat "
            "pertinent (qualité du tout premier bon résultat). nDCG : qualité du classement, avec une "
            "décote logarithmique — un pertinent bien placé compte plus ; normalisé par le classement "
            "idéal (valeur entre 0 et 1).",
            "Calculées dans evaluation/metrics.py sur les 15 requêtes ; nDCG@5 est notre métrique "
            "principale.",
            "Force : objectives, comparables. Limite : dépendent de la vérité terrain.")

    concept("Validation croisée / grid-search & sur-apprentissage",
            "Recherche systématique des meilleurs hyper-paramètres, évaluée sur des données.",
            "On teste toutes les combinaisons de poids (α,β,γ) sommant à 1 et on garde celle qui "
            "maximise le nDCG@5. On retient un compromis (50/40/10) plutôt que l'optimum brut "
            "(50/50/0) pour éviter le sur-apprentissage : l'optimum brut annulait l'ambiance et "
            "collait trop à une vérité terrain bâtie par catégorie.",
            "evaluation/tune_weights.py ; justifie objectivement les poids.",
            "Force : choix fondé sur des données, pas arbitraire. Limite : optimisé sur un petit jeu ; "
            "risque de sur-apprentissage si on suit l'optimum aveuglément.")

    concept("Baseline & comparaison avant/après",
            "Point de référence pour prouver qu'une amélioration apporte réellement de la valeur.",
            "On compare systématiquement : sémantique seul vs pondéré corrigé (preuve du correctif), "
            "et SBERT vs TF-IDF (preuve de l'apport sémantique). Sans baseline, un chiffre isolé ne "
            "prouve rien.",
            "RESULTS.md (avant/après correctif) et baseline_results.md (SBERT vs TF-IDF).",
            "Force : rend les gains démontrables. Limite : la baseline doit être honnête et pertinente.")

    # ------------------------------------------------ D. INGÉNIERIE / PROD
    h1(doc, "D. Ingénierie & mise en production")

    concept("Résilience : retry & backoff exponentiel",
            "Capacité à survivre aux erreurs transitoires d'une API externe.",
            "En cas d'erreur 429 (quota dépassé) ou d'indisponibilité, on relance l'appel après un "
            "délai qui double à chaque tentative (5s, 10s, 20s, 40s…) : c'est le backoff exponentiel, "
            "qui évite de marteler le service et laisse le quota par minute se réinitialiser.",
            "Implémenté dans genai_integration (_generate_with_retry).",
            "Force : robustesse en production. Limite : ne résout pas un quota épuisé pour la journée.")

    concept("Quotas & rate limiting (RPM / RPD)",
            "Limites imposées par l'API : requêtes par minute (RPM) et par jour (RPD).",
            "Le free tier Gemini plafonne le nombre d'appels ; dépasser renvoie une erreur 429. Le "
            "retry gère les limites par minute ; les limites par jour imposent d'attendre la "
            "réinitialisation ou de passer en payant.",
            "Contrainte réelle rencontrée pendant l'évaluation LLM ; gérée par cache + retry + mode "
            "dégradé.",
            "Force : maîtrise des coûts. Limite : bride le volume d'évaluation en free tier.")

    concept("Cache & hachage SHA-256",
            "Mémorisation des réponses pour éviter de rappeler l'API sur une entrée identique.",
            "La clé de cache est un hachage SHA-256 du prompt + modèle + paramètres de génération : "
            "deux réglages différents ne partagent pas la même réponse. Le cache est persistant "
            "(fichier JSON) avec éviction FIFO quand il est plein.",
            "cache_manager.py ; réduit coûts et latence, rend les démos reproductibles.",
            "Force : économies, reproductibilité. Limite : peut servir une réponse périmée si le "
            "prompt ne change pas.")

    concept("Observabilité : suivi tokens, coût, latence",
            "Instrumentation permettant de mesurer ce que consomme le système.",
            "À chaque appel, on lit usage_metadata (tokens entrée/sortie), on estime le coût via une "
            "grille tarifaire, et on chronomètre la latence. On agrège par session.",
            "get_api_stats() ; répond à la question jury « comment industrialiser / maîtriser les "
            "coûts ? ».",
            "Force : pilotage des coûts et de la performance. Limite : coût estimé (grille indicative).")

    concept("Sécurité : injection de prompt & safety settings",
            "Se protéger d'entrées malveillantes et cadrer les contenus générés.",
            "Injection de prompt = un utilisateur glisse « ignore les instructions précédentes… » "
            "dans le texte libre pour détourner le LLM ; on neutralise les lignes suspectes et on "
            "encadre le texte par des balises. Les safety_settings de Gemini bloquent les catégories "
            "à risque (haine, violence…).",
            "sanitize_user_text() + safety_settings dans genai_integration.",
            "Force : gouvernance et sécurité. Limite : filtrage heuristique, non exhaustif.")

    concept("Secrets, .env et .gitignore",
            "Gestion des informations sensibles (clé API) hors du code versionné.",
            "La clé vit dans un fichier .env local, ignoré par git ; le dépôt ne contient qu'un "
            ".env.example avec un placeholder. Une clé avait été commitée par erreur : retirée et "
            "rotée.",
            "Gouvernance (C5.1) ; condition sine qua non d'un projet propre.",
            "Force : pas de fuite de secret. Limite : une clé exposée reste dans l'historique git "
            "(d'où la rotation).")

    concept("Conteneurisation (Docker) & CI (GitHub Actions)",
            "Reproductibilité de l'exécution et automatisation des vérifications.",
            "Docker empaquette l'app et ses dépendances dans une image identique partout. La CI "
            "(intégration continue) exécute automatiquement les tests à chaque push : c'est un "
            "garde-fou de non-régression.",
            "Dockerfile + .github/workflows/ci.yml (29 tests).",
            "Force : déploiement fiable, qualité continue. Limite : image volumineuse (torch) ; CI "
            "limitée aux tests rapides.")

    concept("Tests unitaires & mode dégradé",
            "Vérifier automatiquement des comportements précis ; fonctionner même en panne partielle.",
            "Les tests (pytest) verrouillent le scoring, les métriques, le cache et l'anti-injection. "
            "Le mode dégradé garde le cœur RAG (retrieval + scoring + visualisations) fonctionnel "
            "même sans clé API — seules les synthèses rédigées sont remplacées par un résumé "
            "déterministe.",
            "tests/ + app.py (try/except autour de l'init GenAI).",
            "Force : fiabilité, démo possible sans clé. Limite : couverture de tests à étendre.")

    concept("Streamlit & RGPD",
            "Framework d'UI web en Python ; cadre légal des données personnelles.",
            "Streamlit réexécute le script à chaque interaction (modèle « rerun ») pour construire "
            "l'interface. Les réponses libres de l'utilisateur sont stockées localement et exclues "
            "du dépôt ; en production il faudrait consentement, durée de conservation et anonymisation.",
            "app.py (UI) ; les données utilisateur relèvent du RGPD.",
            "Force : prototypage rapide, accessible. Limite : peu adapté aux très fortes charges.")

    # ------------------------------------------------ E. Q/R + F. mémo
    h1(doc, "E. Questions probables du jury — réponses courtes")
    qa = [
        ("Pourquoi un RAG plutôt qu'un fine-tuning ?",
         "Corpus évolutif sans réentraînement, génération traçable, coût maîtrisé ; le fine-tuning "
         "demanderait beaucoup de données et serait opaque."),
        ("Comment évaluez-vous la qualité ? (LA question C5.3)",
         "Retrieval : 15 requêtes annotées + Precision/Recall/MRR/nDCG. Génération : LLM-as-judge "
         f"sur 5 critères + comparaison de température. Preuve : {NUM['gain_ndcg']} % de nDCG après correctif."),
        ("Pourquoi les poids 50/40/10 ?",
         f"Validés par grid-search : 50/40/10 (nDCG {NUM['w_chosen']}) bat le naïf 50/30/20 "
         f"({NUM['w_naive']}) ; on garde le sémantique dominant et le mood en départage."),
        ("Comment gérez-vous les hallucinations ?",
         "Ancrage RAG, consignes anti-invention, indicateur de titres hors-corpus, critère "
         "« non-hallucination » du juge, température basse possible."),
        ("Et les quotas / la résilience ?",
         "Retry avec backoff exponentiel sur les 429 (limites par jour ET par minute), cache, et "
         "mode dégradé sans clé."),
        ("Qu'avez-vous corrigé exactement ?",
         "Le score de genre comparait du français à une colonne anglaise → 40 % du score inertes ; "
         "on utilise la colonne Categorie (FR) + matching sans accents, verrouillé par des tests. "
         f"Impact : {NUM['gain_ndcg']} % de nDCG."),
        ("Comment industrialiser ?",
         "FAISS + embeddings persistés (faits), Docker + CI (faits), monitoring coûts (instrumenté), "
         "supervision du taux d'hallucination, A/B testing des paramètres."),
        ("Sécurité / RGPD ?",
         "Clé retirée du dépôt et rotée, secrets via .env, safety_settings, anti-injection ; données "
         "utilisateur locales à encadrer (consentement, rétention)."),
    ]
    for q, a in qa:
        para(doc, "Q : " + q, bold=True); para(doc, "R : " + a)

    h1(doc, "F. Mémo transposable aux autres projets")
    for b in ["Relier chaque réalisation à une compétence (« ceci démontre Cx.y car… »).",
              "Suivre la logique besoin → choix → réalisation → preuve → résultat → limites.",
              "Toujours chiffrer les résultats et comparer à une baseline (avant/après).",
              "Assumer les limites et proposer des améliorations (le jury valorise le recul).",
              "Mentionner systématiquement sécurité, reproductibilité (deps figées, .env, Docker), tests, CI.",
              "Savoir expliquer chaque brique simplement, puis en profondeur si on creuse (cf. sections A-D)."]:
        bullet(doc, b)

    doc.save(OUT / "antiseche_projet3.docx")
    print("OK antiseche_projet3.docx")


# ======================================================= 4. DOSSIER AUDIT
def build_dossier():
    doc = Document()
    title_block(doc, "Dossier d'audit & finalisation — Projet 3",
                "Audit · Analyse d'écart RNCP · Récapitulatif final")

    h1(doc, "Partie 1 — Audit de l'existant")
    para(doc, "Base technique saine (RAG réel, code modulaire, cache) mais plusieurs problèmes "
              "identifiés et prouvés par exécution :")
    bullet(doc, "BUG MAJEUR : composant genre (30 % du score) neutralisé — libellés FR vs colonne Genre anglaise.")
    bullet(doc, "Mood matché seulement par hasard (divergences d'accents).")
    bullet(doc, "Modèle Gemini incohérent ; ré-encodage du référentiel à chaque requête.")
    bullet(doc, "Aucune évaluation (C5.3) ; génération sans paramétrage ni résilience.")
    bullet(doc, "Clé API réelle commitée dans .env.example.")

    h1(doc, "Partie 2 — Analyse d'écart vs grille RNCP")
    para(doc, "Échelle : Non démontré (0) / Débutant (2) / Intermédiaire (3) / Professionnel (5).")
    table(doc, ["Compétence", "Avant", "Manque principal", "Action corrective"], [
        ["C5.1 Cas d'usage", "3", "besoin métier non formalisé, gouvernance", "rapport, retrait/rotation clé, safety"],
        ["C5.2 Solution", "3", "bug genre bloquant, traçabilité, perf, résilience", "fix scoring, sources UI, FAISS, retry"],
        ["C5.3 Évaluation", "0", "aucune évaluation", "métriques + calibrage poids + LLM-as-judge + compare params"],
    ])

    h1(doc, "Partie 3 — Récapitulatif final")
    table(doc, ["Compétence", "Avant", "Après", "Preuve"], [
        ["C5.1", "3/5", "5/5", "rapport ; .env.example assaini + rotation ; safety_settings ; anti-injection"],
        ["C5.2", "3/5", "5/5", "RAG fonctionnel/résilient, UI + traçabilité, fix scoring (tests), FAISS"],
        ["C5.3", "0/5", "5/5", f"métriques (+{NUM['gain_ndcg'].lstrip('+')}% nDCG), calibrage poids, LLM-as-judge, compare params"],
    ])
    h2(doc, "Actions réalisées")
    for a in ["Clé API retirée de .env.example + rotation",
              "Scoring genre/mood corrigé (Categorie FR + accents) + tests",
              "Calibrage des poids par validation croisée (50/40/10)",
              "Cadre d'évaluation : métriques retrieval + baseline TF-IDF + LLM-as-judge",
              "GenerationConfig ajustable + comparaison de paramètres (température)",
              "Résilience : retry/backoff sur quota 429, mode dégradé sans clé",
              "Gouvernance : safety_settings, garde-fou anti-prompt-injection",
              "Observabilité : suivi tokens / coût estimé / latence",
              "Perf : persistance des embeddings + index FAISS optionnel",
              "Qualité : 29 tests unitaires + CI GitHub Actions",
              "Conteneurisation : Dockerfile",
              "Livrables : rapport, slides, plan, antisèche, dossier"]:
        bullet(doc, "Fait — " + a)
    h2(doc, "Limites restantes assumées")
    bullet(doc, "Corpus limité (260 films EN) ; vérité terrain par les auteurs.")
    bullet(doc, "Hallucination résiduelle hors top 3 ; quotas free tier (jour/minute).")
    bullet(doc, "CI limitée aux tests (l'évaluation complète reste lançable en local).")

    doc.save(OUT / "dossier_audit_projet3.docx")
    print("OK dossier_audit_projet3.docx")


# ======================================================= 5. PRÉSENTATION PPTX
# Système de design repris du modèle fourni (deck "Urban Data Explorer") :
# fond clair, badges numérotés verts, pills, cartes arrondies, puces carrées.
INK = PRGB(0x14, 0x28, 0x2E)      # encre (titres/texte)
GREEN = PRGB(0x24, 0x6B, 0x63)    # vert primaire (badges, accents)
GREEN2 = PRGB(0x3F, 0x8F, 0x84)   # vert secondaire
MUTED = PRGB(0x5A, 0x70, 0x77)    # gris-vert (texte secondaire)
ORANGE = PRGB(0xD2, 0x6A, 0x4E)   # accent orange (chiffres clés)
BG = PRGB(0xF3, 0xF6, 0xF5)       # fond
WHITE = PRGB(0xFF, 0xFF, 0xFF)
CARD_LT = PRGB(0xDC, 0xED, 0xE9)  # carte / pill claire
LIGHT = PRGB(0xE9, 0xF1, 0xEF)
FT_TITLE = "Cambria"
FT_BODY = "Calibri"

IN = PInches


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _round(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = PPt(1)
    sp.shadow.inherit = False
    return sp


def _sq(slide, l, t, fill, size=0.16):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l), IN(t), IN(size), IN(size))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _txt(slide, l, t, w, h, text, size, color, font=FT_BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = PPt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tb


def _header(slide, num, title, tag=None):
    """Bandeau : badge numéroté + titre + pill (comme le modèle)."""
    b = _round(slide, 0.6, 0.55, 0.62, 0.62, GREEN)
    tf = b.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.name = FT_TITLE; r.font.size = PPt(24); r.font.bold = True; r.font.color.rgb = WHITE
    _txt(slide, 1.4, 0.5, 8.6, 0.75, title, 27, INK, font=FT_TITLE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    if tag:
        pill = _round(slide, 10.05, 0.66, 2.68, 0.5, CARD_LT)
        tfp = pill.text_frame; tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tfp.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rp = pp.add_run(); rp.text = tag
        rp.font.name = FT_BODY; rp.font.size = PPt(12); rp.font.bold = True; rp.font.color.rgb = GREEN
    # filet vert sous le bandeau
    _round(slide, 0.6, 1.4, 12.13, 0.03, GREEN)


def _bullets(slide, items, top=1.9, x=0.9, w=11.6, gap=1.02, sq=ORANGE):
    for i, it in enumerate(items):
        y = top + i * gap
        _sq(slide, x, y + 0.06, sq)
        if isinstance(it, (tuple, list)):
            _txt(slide, x + 0.4, y - 0.05, w, 0.45, it[0], 18, INK, bold=True)
            _txt(slide, x + 0.4, y + 0.34, w, 0.55, it[1], 14, MUTED)
        else:
            _txt(slide, x + 0.4, y - 0.02, w, 0.7, it, 17, INK)


def _stat_cards(slide, stats, top=4.4):
    """Rangée de cartes chiffres (label + valeur en accent)."""
    n = len(stats); gapx = 0.3
    w = (12.13 - gapx * (n - 1)) / n
    for i, (val, label, color) in enumerate(stats):
        x = 0.6 + i * (w + gapx)
        _round(slide, x, top, w, 1.6, WHITE, line=CARD_LT)
        _txt(slide, x, top + 0.22, w, 0.7, val, 32, color, font=FT_TITLE, bold=True,
             align=PP_ALIGN.CENTER)
        _txt(slide, x, top + 1.02, w, 0.5, label, 12, MUTED, align=PP_ALIGN.CENTER)


def _image_card(slide, img, l, t, w, h):
    if img and Path(img).exists():
        _round(slide, l - 0.12, t - 0.12, w + 0.24, h + 0.24, WHITE, line=CARD_LT)
        slide.shapes.add_picture(str(img), IN(l), IN(t), height=IN(h))


def _new(prs, layout):
    s = prs.slides.add_slide(layout); _bg(s); return s


def build_pptx():
    # Deck vierge : le système de design ci-dessus applique explicitement la
    # palette et les polices du modèle fourni (Urban Data Explorer), donc le
    # rendu reprend le même style sans dépendre du fichier source (évite aussi
    # les parties orphelines d'un template dont on retirerait les slides).
    prs = Presentation()
    layout = prs.slide_layouts[6]
    prs.slide_width = IN(13.333); prs.slide_height = IN(7.5)

    # --- Slide 1 : titre ---
    s = _new(prs, layout)
    _round(s, 0.9, 2.15, 0.18, 1.9, GREEN)  # barre d'accent verticale
    _txt(s, 1.3, 2.05, 11, 1.3, "AISCA-Cinema", 52, INK, font=FT_TITLE, bold=True)
    _txt(s, 1.32, 3.35, 11.4, 0.8,
         "Agent de recommandation cinématographique (RAG)", 24, GREEN, font=FT_TITLE)
    pill = _round(s, 1.32, 4.5, 4.6, 0.55, CARD_LT)
    tfp = pill.text_frame; tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tfp.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rp = pp.add_run(); rp.text = "Projet 3 · IA Générative · Bloc 2"
    rp.font.name = FT_BODY; rp.font.size = PPt(13); rp.font.bold = True; rp.font.color.rgb = GREEN
    _txt(s, 1.32, 5.5, 11, 0.5, AUTHORS, 20, INK, bold=True)
    _txt(s, 1.32, 6.0, 11, 0.5, "RNCP40875 — Expert en Ingénierie des données (C5.1 → C5.3)",
         14, MUTED)

    # --- Slide 2 : besoin métier ---
    s = _new(prs, layout)
    _header(s, 2, "Besoin métier", "🎯 C5.1")
    _bullets(s, [
        ("Trop de films, choix coûteux", "Les filtres par mots-clés ne comprennent pas une envie en langage naturel."),
        ("Cible", "Un spectateur qui décrit une envie… mais pas un titre précis."),
        ("Valeur", "Recommander ET expliquer, à partir d'un référentiel maîtrisé."),
        ("Pourquoi la GenAI", "Compréhension sémantique + restitution personnalisée, hors de portée d'un simple filtre."),
    ])

    # --- Slide 3 : architecture RAG ---
    s = _new(prs, layout)
    _header(s, 3, "Architecture RAG", "🧭 Vue d'ensemble")
    steps = ["Questionnaire", "SBERT", "Scoring", "Gemini"]
    x = 0.9; cw = 2.5; gap = 0.55
    for i, st in enumerate(steps):
        cx = x + i * (cw + gap)
        c = _round(s, cx, 1.95, cw, 0.95, GREEN if i in (1, 3) else CARD_LT)
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = st; r.font.name = FT_BODY; r.font.size = PPt(15)
        r.font.bold = True; r.font.color.rgb = WHITE if i in (1, 3) else INK
        if i < len(steps) - 1:
            _txt(s, cx + cw, 1.95, gap, 0.95, "→", 22, GREEN, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, 0.9, 3.05, 11.6, 0.4, "Retrieval (récupération sémantique)  →  Generation (synthèse ancrée)",
         13, MUTED)
    _image_card(s, SHOTS / "01_questionnaire.png", 4.55, 3.7, 4.2, 3.3)
    _txt(s, 0.9, 3.7, 3.4, 3.2,
         "• Embeddings SBERT multilingue + cosinus (260 films)\n\n"
         "• Score 0.50 sém. + 0.40 genre + 0.10 mood\n\n"
         "• Génération Gemini ancrée sur les sources\n\n"
         "• Interface Streamlit accessible", 14, INK)

    # --- Slide 4 : choix techniques ---
    s = _new(prs, layout)
    _header(s, 4, "Choix techniques", "⚙️ C5.2")
    _bullets(s, [
        ("RAG", "Corpus évolutif sans réentraînement, génération traçable, coût maîtrisé."),
        ("vs Fine-tuning / Prompting seul", "Écartés : coût & opacité / risque d'hallucination sans ancrage."),
        (f"Embeddings vs mots-clés : SBERT +{NUM['tfidf_gain']} %",
         "Baseline TF-IDF s'effondre (requêtes FR / films EN) : seul le multilingue fait le pont."),
        ("Coûts maîtrisés", "Cache + retry/backoff + 2 appels API par session."),
    ])

    # --- Slide 5 : démo Top 3 ---
    s = _new(prs, layout)
    _header(s, 5, "Démo — Top 3 recommandations", "▶️ Live")
    _image_card(s, SHOTS / "02_top3.png", 3.4, 1.75, 6.5, 5.1)
    _txt(s, 0.9, 6.75, 11.6, 0.5,
         "Sortie réelle : Top 3 + décomposition du score (sémantique / genre / mood).",
         13, MUTED, align=PP_ALIGN.CENTER)

    # --- Slide 6 : visualisations ---
    s = _new(prs, layout)
    _header(s, 6, "Démo — Visualisations", "📊 UX")
    _image_card(s, SHOTS / "03_visualisations.png", 3.4, 1.75, 6.5, 5.1)
    _txt(s, 0.9, 6.75, 11.6, 0.5,
         "Profils par genre et ambiance ; score pondéré 50/40/10 par film.",
         13, MUTED, align=PP_ALIGN.CENTER)

    # --- Slide 7 : évaluation (chiffres) ---
    s = _new(prs, layout)
    _header(s, 7, "Évaluation — avant / après", "📈 C5.3")
    _image_card(s, FIGS / "rag_comparison.png", 0.7, 1.75, 6.6, 3.9)
    _txt(s, 7.6, 1.9, 5.2, 2.6,
         "• 15 requêtes annotées (vérité terrain)\n\n"
         "• Barre orange (genre buggé) ≈ sémantique : preuve du bug\n\n"
         "• Poids 50/40/10 validés par grid-search\n\n"
         "• LLM-as-judge + comparaison de température", 14, INK)
    _stat_cards(s, [
        (f"{NUM['ndcg_sem']}→{NUM['ndcg_cal']}", "nDCG@5", GREEN),
        (f"{NUM['gain_ndcg']} %", "gain nDCG", ORANGE),
        (f"{NUM['mrr_sem']}→{NUM['mrr_cal']}", "MRR", GREEN),
        (f"+{NUM['tfidf_gain'].lstrip('+')} %", "vs TF-IDF", ORANGE),
    ], top=5.75)

    # --- Slide 8 : robustesse / sécurité / limites ---
    s = _new(prs, layout)
    _header(s, 8, "Robustesse · Sécurité · Limites", "🛡️ Prod")
    _bullets(s, [
        ("Résilience", "Retry/backoff sur quota 429 + mode dégradé sans clé."),
        ("Gouvernance", "safety_settings + garde-fou anti-prompt-injection."),
        ("Observabilité", "Suivi tokens / coût estimé / latence par appel."),
        ("Limites assumées", "Corpus 260 films EN, hallucination résiduelle, quotas free tier."),
    ])

    # --- Slide 9 : industrialisation + conclusion ---
    s = _new(prs, layout)
    _header(s, 9, "Industrialisation & conclusion", "🚀 Bilan")
    _bullets(s, [
        ("Industrialisation", "FAISS + embeddings persistés, Docker + CI, monitoring coûts, A/B testing."),
        ("C5.1", "Cas d'usage justifié et gouverné."),
        ("C5.2", "Solution RAG fonctionnelle, accessible, résiliente."),
        (f"C5.3", f"Évaluation chiffrée avec optimisation prouvée ({NUM['gain_ndcg']} % nDCG)."),
    ], top=1.85, gap=0.92)
    _txt(s, 0.9, 6.6, 11.6, 0.5, "Merci — questions ?", 20, GREEN, font=FT_TITLE,
         bold=True, align=PP_ALIGN.CENTER)

    prs.save(OUT / "presentation_projet3.pptx")
    print("OK presentation_projet3.pptx")


if __name__ == "__main__":
    _load_numbers()
    build_rapport()
    build_plan()
    build_antiseche()
    build_dossier()
    build_pptx()
    print("\nLivrables dans", OUT)
